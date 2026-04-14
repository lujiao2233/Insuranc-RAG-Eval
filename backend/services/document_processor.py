"""文档处理器模块
处理多种格式文档的解析和文本提取功能
支持PDF、DOCX、XLSX格式
"""
import os
from pathlib import Path
from typing import Optional, Dict, Any
from datetime import datetime
import PyPDF2
from docx import Document as DocxDocument
import openpyxl
from PIL import Image
import io

from utils.logger import get_logger
from services.metadata_extractor import MetadataExtractor

logger = get_logger("document_processor")


class DocumentProcessor:
    """文档处理服务类
    
    提供文档解析、文本提取和元数据提取的完整功能。
    支持PDF、DOCX、XLSX、PPTX、TXT格式，并保留结构标记。
    """
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self.process_pdf,
            'docx': self.process_docx,
            'xlsx': self.process_xlsx,
            'pptx': self.process_pptx,
            'txt': self.process_txt,
            'excel': self.process_xlsx  # 别名
        }
    
    def process_file(self, file_path: str) -> Optional[Dict[str, Any]]:
        """处理文档文件"""
        try:
            path = Path(file_path)
            if not path.exists():
                logger.error(f"文件不存在: {file_path}")
                return None
            
            file_extension = path.suffix.lower().lstrip('.')
            if file_extension not in self.supported_formats:
                logger.error(f"不支持的文件格式: {file_extension}")
                return None
            
            file_size = path.stat().st_size
            
            processor = self.supported_formats[file_extension]
            result = processor(str(path))
            
            if result is None:
                return None
            
            text_content = result.get("text_content", "")
            
            return {
                "filename": path.name,
                "file_path": str(path),
                "file_type": file_extension,
                "file_size": file_size,
                "page_count": result.get("page_count"),
                "text_content": text_content,
                "metadata": result.get("metadata", {}),
                "outline": result.get("outline", [])
            }
            
        except Exception as e:
            logger.error(f"文档处理失败: {str(e)}")
            return None

    def process_txt(self, file_path: str) -> Optional[Dict[str, Any]]:
        """处理TXT文件"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return {
                "page_count": 1,
                "text_content": content,
                "metadata": {"paragraph_count": len(content.split('\n\n'))}
            }
        except Exception as e:
            logger.error(f"TXT处理失败: {str(e)}")
            return None

    def process_pptx(self, file_path: str) -> Optional[Dict[str, Any]]:
        """处理PPTX文件"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = ""
            for i, slide in enumerate(prs.slides):
                text_content += f"\n[SLIDE {i+1}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += f"[PARAGRAPH] {shape.text}\n"
                    # 处理图片占位符（如果PPT中有alt文字，可以在此提取）
                    if shape.shape_type == 13: # Picture
                        text_content += f"[IMAGE_ALT: {shape.name}]\n"
            return {
                "page_count": len(prs.slides),
                "text_content": text_content.strip(),
                "metadata": {"slide_count": len(prs.slides)}
            }
        except ImportError:
            logger.warning("python-pptx未安装，PPTX处理失败")
            return None
        except Exception as e:
            logger.error(f"PPTX处理失败: {str(e)}")
            return None

    def process_pdf(self, file_path: str) -> Optional[Dict[str, Any]]:
        """处理PDF文件，保留结构标记"""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            text_content = ""
            for page in doc:
                text_content += f"\n[PAGE {page.number + 1}]\n"
                # 尝试提取结构化块
                blocks = page.get_text("blocks")
                for b in blocks:
                    # b: (x0, y0, x1, y1, "text", block_no, block_type)
                    if b[6] == 0:  # Text block
                        text_content += f"[PARAGRAPH] {b[4]}\n"
                    elif b[6] == 1:  # Image block
                        text_content += "[IMAGE_ALT: PDF Image]\n"
                
                # 尝试提取表格
                try:
                    tabs = page.find_tables()
                    for i, tab in enumerate(tabs):
                        text_content += f"\n[TABLE {i+1}]\n"
                        for row in tab.extract():
                            text_content += " | ".join([str(c) if c else "" for c in row]) + "\n"
                except:
                    pass
            doc.close()
            return {
                "page_count": len(doc),
                "text_content": text_content.strip(),
                "metadata": {}
            }
        except ImportError:
            logger.info("PyMuPDF未安装，回退到PyPDF2")
            import pypdf
            reader = pypdf.PdfReader(file_path)
            text_content = ""
            for i, page in enumerate(reader.pages):
                text_content += f"\n[PAGE {i+1}]\n"
                text_content += f"[PARAGRAPH] {page.extract_text()}\n"
            return {
                "page_count": len(reader.pages),
                "text_content": text_content.strip(),
                "metadata": {}
            }
        except Exception as e:
            logger.error(f"PDF处理失败: {str(e)}")
            return None

    def process_docx(self, file_path: str) -> Optional[Dict[str, Any]]:
        """处理DOCX文件，保留结构标记"""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            text_content = ""
            
            # 遍历文档的所有元素，保持顺序
            for element in doc.element.body:
                if element.tag.endswith('p'): # Paragraph
                    from docx.text.paragraph import Paragraph
                    p = Paragraph(element, doc)
                    if p.text.strip():
                        text_content += f"[PARAGRAPH] {p.text}\n"
                elif element.tag.endswith('tbl'): # Table
                    from docx.table import Table
                    t = Table(element, doc)
                    text_content += "\n[TABLE]\n"
                    for row in t.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        text_content += row_text + "\n"
            
            return {
                "page_count": max(1, len(doc.paragraphs) // 20),
                "text_content": text_content.strip(),
                "metadata": {"paragraph_count": len(doc.paragraphs)}
            }
        except Exception as e:
            logger.error(f"DOCX处理失败: {str(e)}")
            return None
    
    def read_document_content(self, file_path: str) -> Optional[str]:
        """读取文档内容
        
        仅提取文档的文本内容，不进行元数据提取。
        """
        try:
            path = Path(file_path)
            if not path.exists():
                logger.error(f"文件不存在: {file_path}")
                return None
            
            file_extension = path.suffix.lower().lstrip('.')
            if file_extension not in self.supported_formats:
                logger.error(f"不支持的文件格式: {file_extension}")
                return None
            
            processor = self.supported_formats[file_extension]
            result = processor(str(path))
            
            if result:
                return result.get("text_content", "")
            return None
            
        except Exception as e:
            logger.error(f"读取文档内容失败: {str(e)}")
            return None
    
    def _extract_table_like_blocks(self, text: str) -> Dict[str, Any]:
        """提取表格类文本块
        
        通过分析文本结构识别可能的表格内容。
        """
        lines = text.splitlines()
        blocks = []
        current_block = []
        
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if current_block:
                    blocks.append("\n".join(current_block))
                    current_block = []
                continue
            
            has_tab = "\t" in stripped
            parts = [p for p in stripped.split(" ") if p]
            has_many_columns = len(parts) >= 3
            has_pipe = "|" in stripped
            
            if has_tab or has_many_columns or has_pipe:
                current_block.append(stripped)
            else:
                if current_block:
                    blocks.append("\n".join(current_block))
                    current_block = []
        
        if current_block:
            blocks.append("\n".join(current_block))
        
        return {
            "table_like_blocks": blocks,
            "table_like_count": len(blocks),
        }

    def process_xlsx(self, file_path: str) -> Optional[Dict[str, Any]]:
        """处理XLSX文件
        
        使用openpyxl库解析Excel电子表格，提取工作表信息和数据内容。
        """
        try:
            if not os.path.exists(file_path):
                logger.warning(f"XLSX文件不存在: {file_path}")
                return None
            
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            
            sheet_count = len(workbook.sheetnames)
            
            first_sheet = workbook[workbook.sheetnames[0]]
            row_count = first_sheet.max_row
            column_count = first_sheet.max_column
            
            sample_data = []
            for row in first_sheet.iter_rows(min_row=1, max_row=row_count, 
                                           values_only=True):
                sample_data.append([str(cell) if cell is not None else "" 
                                  for cell in row])
            
            sheet_names = list(workbook.sheetnames)
            workbook.close()
            
            return {
                "page_count": sheet_count,
                "text_content": str(sample_data),
                "metadata": {
                    "sheet_count": sheet_count,
                    "row_count": row_count,
                    "column_count": column_count,
                    "sheet_names": sheet_names
                }
            }
            
        except Exception as e:
            logger.error(f"XLSX处理失败: {str(e)}")
            return None
    
    def get_document_preview(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """获取文档预览信息
        
        生成文档的预览信息字典，用于在UI中展示文档概要。
        """
        try:
            path = Path(file_path)
            if not path.exists():
                return {"error": "文件不存在"}
            
            stat = path.stat()
            
            preview_info = {
                "filename": path.name,
                "file_size": stat.st_size,
                "file_size_mb": round(stat.st_size / (1024 * 1024), 2),
                "file_type": file_type,
                "modified_time": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
            }
            
            if file_type == "pdf":
                pdf_info = self.process_pdf(str(path))
                if pdf_info:
                    preview_info.update({
                        "page_count": pdf_info.get("page_count", 0),
                        "title": pdf_info["metadata"].get("title", ""),
                        "author": pdf_info["metadata"].get("author", "")
                    })
            
            elif file_type == "docx":
                docx_info = self.process_docx(str(path))
                if docx_info:
                    preview_info.update({
                        "page_count": docx_info.get("page_count", 0),
                        "paragraph_count": docx_info["metadata"].get("paragraph_count", 0),
                        "table_count": docx_info["metadata"].get("table_count", 0)
                    })
            
            elif file_type == "xlsx":
                xlsx_info = self.process_xlsx(str(path))
                if xlsx_info:
                    preview_info.update({
                        "page_count": xlsx_info.get("page_count", 0),
                        "sheet_count": xlsx_info["metadata"].get("sheet_count", 0),
                        "sheet_names": xlsx_info["metadata"].get("sheet_names", [])
                    })
            
            return preview_info
            
        except Exception as e:
            logger.error(f"获取文档预览信息失败: {str(e)}")
            return {"error": str(e)}
